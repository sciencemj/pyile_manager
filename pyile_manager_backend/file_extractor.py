"""
File content extraction utilities for different file types.
Supports PDF, images, PowerPoint, and text files.
"""

import base64
from pathlib import Path

import pdfplumber
from PIL import Image
from pptx import Presentation


def extract_text_from_pdf(file_path: str) -> str | None:
    """
    Extract text content from a PDF file.
    Returns None if extraction fails or if PDF has no text (scanned document).
    """
    try:
        text_content = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_content.append(page_text)

        full_text = "\n".join(text_content).strip()
        # If text is too short, it's likely a scanned PDF
        if len(full_text) < 50:
            return None
        return full_text
    except Exception as e:
        print(f"Error extracting text from PDF {file_path}: {e}")
        return None


def extract_text_from_pptx(file_path: str) -> str | None:
    """
    Extract text content from a PowerPoint file.
    """
    try:
        prs = Presentation(file_path)
        text_content = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)

        return "\n".join(text_content).strip()
    except Exception as e:
        print(f"Error extracting text from PPTX {file_path}: {e}")
        return None


def extract_text_from_txt(file_path: str) -> str | None:
    """
    Extract text content from a plain text file.
    """
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except UnicodeDecodeError:
        # Try with different encoding
        try:
            with open(file_path, "r", encoding="latin-1") as f:
                return f.read().strip()
        except Exception as e:
            print(f"Error reading text file {file_path}: {e}")
            return None
    except Exception as e:
        print(f"Error reading text file {file_path}: {e}")
        return None


def is_image_file(file_path: str) -> bool:
    """
    Check if file is an image based on extension.
    """
    image_extensions = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff"}
    return Path(file_path).suffix.lower() in image_extensions


def is_pdf_file(file_path: str) -> bool:
    """
    Check if file is a PDF.
    """
    return Path(file_path).suffix.lower() == ".pdf"


def is_pptx_file(file_path: str) -> bool:
    """
    Check if file is a PowerPoint file.
    """
    ppt_extensions = {".ppt", ".pptx"}
    return Path(file_path).suffix.lower() in ppt_extensions


def is_text_file(file_path: str) -> bool:
    """
    Check if file is a plain text file.
    """
    return Path(file_path).suffix.lower() == ".txt"


def get_file_type(file_path: str) -> str:
    """
    Determine the file type category.
    Returns: 'image', 'pdf', 'pptx', 'text', or 'unknown'
    """
    if is_image_file(file_path):
        return "image"
    elif is_pdf_file(file_path):
        return "pdf"
    elif is_pptx_file(file_path):
        return "pptx"
    elif is_text_file(file_path):
        return "text"
    else:
        return "unknown"
